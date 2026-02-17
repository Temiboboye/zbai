#!/usr/bin/env python3
"""
Convert ZeroBounce AI HTML Pitch Deck to PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
GREEN = RGBColor(127, 217, 87)  # #7FD957
DARK = RGBColor(10, 10, 10)  # #0A0A0A
WHITE = RGBColor(255, 255, 255)
GRAY_600 = RGBColor(75, 85, 99)
GRAY_700 = RGBColor(55, 65, 81)

def add_title_slide(prs):
    """Slide 1: Title"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Logo
    logo_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(4), Inches(0.8))
    logo_frame = logo_box.text_frame
    logo_frame.text = "ZB ZeroBounce AI"
    logo_p = logo_frame.paragraphs[0]
    logo_p.alignment = PP_ALIGN.CENTER
    logo_run = logo_p.runs[0]
    logo_run.font.size = Pt(36)
    logo_run.font.bold = True
    logo_run.font.color.rgb = DARK
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Email Verification That\nActually Thinks"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.runs[0]
    title_run.font.size = Pt(54)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Email Verification with 98%+ Accuracy"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.size = Pt(24)
    subtitle_run.font.color.rgb = GRAY_600
    
    # Founder info
    founder_box = slide.shapes.add_textbox(Inches(2.5), Inches(5.5), Inches(5), Inches(0.8))
    founder_frame = founder_box.text_frame
    founder_frame.text = "Temitayo Boboye - Founder\nhello@zerobounceai.com"
    for p in founder_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = DARK

def add_problem_slide(prs):
    """Slide 2: Problem"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Email Verification Is Broken"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Problem 1
    prob1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(2.8), Inches(2))
    prob1_frame = prob1_box.text_frame
    prob1_frame.text = "❌ Traditional Tools: 92-95% Accuracy\n\n• 5-8% of emails still bounce\n• Damaged sender reputation\n• Lost revenue opportunities"
    for p in prob1_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_700
    
    # Problem 2
    prob2_box = slide.shapes.add_textbox(Inches(3.6), Inches(1.8), Inches(2.8), Inches(2))
    prob2_frame = prob2_box.text_frame
    prob2_frame.text = "❌ Catch-All Domains = Guesswork\n\n• Binary yes/no answers\n• No confidence scoring\n• Marketers waste 30% of budget"
    for p in prob2_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_700
    
    # Problem 3
    prob3_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.8), Inches(2.8), Inches(2))
    prob3_frame = prob3_box.text_frame
    prob3_frame.text = "❌ No Pattern Recognition\n\n• Can't find correct email formats\n• Manual research required\n• Slow, inefficient processes"
    for p in prob3_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_700
    
    # Cost highlight
    cost_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1))
    cost_frame = cost_box.text_frame
    cost_frame.text = "💰 $2.5B wasted annually on bounced emails and poor targeting"
    cost_p = cost_frame.paragraphs[0]
    cost_p.alignment = PP_ALIGN.CENTER
    cost_run = cost_p.runs[0]
    cost_run.font.size = Pt(22)
    cost_run.font.bold = True
    cost_run.font.color.rgb = WHITE
    
    # Add colored background for cost highlight
    cost_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(5), Inches(7), Inches(1)
    )
    cost_shape.fill.solid()
    cost_shape.fill.fore_color.rgb = RGBColor(255, 107, 107)
    cost_shape.line.color.rgb = RGBColor(255, 107, 107)
    
    # Move text box to front
    slide.shapes._spTree.remove(cost_box._element)
    slide.shapes._spTree.append(cost_box._element)

def add_solution_slide(prs):
    """Slide 3: Solution"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "AI-Powered Verification with 98%+ Accuracy"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Solution boxes
    solutions = [
        ("✅ AI-Powered Verification", "• 98%+ accuracy (vs 92-95%)\n• Reduces bounces by 60%\n• Machine learning models"),
        ("✅ Catch-All Confidence Scoring", "• 0-100 confidence score\n• Not just 'maybe' - know exactly\n• Make data-driven decisions"),
        ("✅ Email Pattern Recognition", "• AI suggests correct formats\n• Finds emails competitors miss\n• Saves hours of manual work"),
        ("✅ Domain Reputation Intelligence", "• Protects your sender score\n• Identifies risky domains\n• Prevents blacklisting")
    ]
    
    positions = [(0.5, 1.8), (5, 1.8), (0.5, 4.2), (5, 4.2)]
    
    for (title, content), (x, y) in zip(solutions, positions):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.2), Inches(2))
        frame = box.text_frame
        frame.text = f"{title}\n\n{content}"
        
        # Style title
        title_p = frame.paragraphs[0]
        title_run = title_p.runs[0]
        title_run.font.size = Pt(16)
        title_run.font.bold = True
        title_run.font.color.rgb = DARK
        
        # Style content
        for p in frame.paragraphs[1:]:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.color.rgb = GRAY_700

def add_demo_slide(prs):
    """Slide 4: Product Demo"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "See It In Action"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Input
    input_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(3), Inches(1))
    input_frame = input_box.text_frame
    input_frame.text = "INPUT\n\njohn@example.com"
    input_p = input_frame.paragraphs[0]
    input_run = input_p.runs[0]
    input_run.font.size = Pt(14)
    input_run.font.bold = True
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(1), Inches(0.5))
    arrow_frame = arrow_box.text_frame
    arrow_frame.text = "→"
    arrow_p = arrow_frame.paragraphs[0]
    arrow_p.alignment = PP_ALIGN.CENTER
    arrow_run = arrow_p.runs[0]
    arrow_run.font.size = Pt(36)
    arrow_run.font.color.rgb = GREEN
    
    # Output
    output_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(3.5))
    output_frame = output_box.text_frame
    output_frame.text = """OUTPUT

✅ Status: Valid (Safe)
Confidence: 98/100
✅ MX Records: Found
✅ SMTP: Responsive
⚠️ Catch-All: Yes (High Conf.)
🟢 Domain Rep: Excellent (95)
Suggested: john.doe@example.com"""
    
    for p in output_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = DARK
    
    # Stats
    stats = [
        ("< 200ms", "Response Time"),
        ("98.3%", "Accuracy Rate"),
        ("10 min", "Integration Time")
    ]
    
    x_positions = [1.5, 4, 6.5]
    for (value, label), x in zip(stats, x_positions):
        stat_box = slide.shapes.add_textbox(Inches(x), Inches(5.5), Inches(2), Inches(1))
        stat_frame = stat_box.text_frame
        stat_frame.text = f"{value}\n{label}"
        
        # Value
        value_p = stat_frame.paragraphs[0]
        value_p.alignment = PP_ALIGN.CENTER
        value_run = value_p.runs[0]
        value_run.font.size = Pt(28)
        value_run.font.bold = True
        value_run.font.color.rgb = GREEN
        
        # Label
        label_p = stat_frame.paragraphs[1]
        label_p.alignment = PP_ALIGN.CENTER
        for run in label_p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_600

def add_market_slide(prs):
    """Slide 5: Market Opportunity"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Massive Market, Growing Fast"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # TAM/SAM/SOM
    markets = [
        ("TAM", "$2.5B", "Email Verification", RGBColor(102, 126, 234)),
        ("SAM", "$850M", "AI-Powered Segment", RGBColor(240, 147, 251)),
        ("SOM", "$85M", "3-Year Target", GREEN)
    ]
    
    y_pos = 1.8
    for label, value, desc, color in markets:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(4), Inches(1.2))
        frame = box.text_frame
        frame.text = f"{label}\n{value}\n{desc}"
        
        # Style
        for i, p in enumerate(frame.paragraphs):
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                if i == 0:  # Label
                    run.font.size = Pt(14)
                    run.font.bold = True
                elif i == 1:  # Value
                    run.font.size = Pt(36)
                    run.font.bold = True
                else:  # Description
                    run.font.size = Pt(14)
                run.font.color.rgb = DARK
        
        y_pos += 1.5
    
    # Market trends
    trends_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(2.5))
    trends_frame = trends_box.text_frame
    trends_frame.text = """📈 Market Trends

• Email marketing: 20% YoY growth
• AI adoption: 35% CAGR
• GDPR/compliance: Increasing demand
• API-first tools: 40% YoY growth"""
    
    for p in trends_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_700
    
    # Target customers
    customers_box = slide.shapes.add_textbox(Inches(5.5), Inches(4.5), Inches(4), Inches(1.8))
    customers_frame = customers_box.text_frame
    customers_frame.text = """🎯 Target Customers

SaaS Companies | Marketing Agencies
E-commerce Platforms | Sales Automation Tools"""
    
    for p in customers_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = GRAY_700

def add_business_model_slide(prs):
    """Slide 6: Business Model"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Simple, Scalable, Profitable"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Revenue model
    revenue_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2))
    revenue_frame = revenue_box.text_frame
    revenue_frame.text = """💰 Revenue Model
Credit-based SaaS

• 1 credit = 1 email verification
• Credits rollover (6-24 months)
• No subscriptions - pay as you grow"""
    
    for p in revenue_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_700
    
    # Pricing
    pricing_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(2))
    pricing_frame = pricing_box.text_frame
    pricing_frame.text = """📦 Pricing Tiers

Starter: 1K credits - $8.40
Pro: 10K credits - $55
Business: 50K credits - $209
Enterprise: Custom pricing"""
    
    for p in pricing_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY_700
    
    # Unit economics
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
    metrics_frame = metrics_box.text_frame
    metrics_frame.text = """📊 Unit Economics

$127 ARPU/month  |  $45 CAC  |  $1,524 LTV  |  34x LTV/CAC  |  87% Gross Margin  |  <1 mo Payback"""
    
    for p in metrics_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = DARK

def add_traction_slide(prs):
    """Slide 7: Traction"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Strong Growth, Proven Product-Market Fit"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(38)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Growth chart placeholder
    chart_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
    chart_frame = chart_box.text_frame
    chart_frame.text = """📈 MRR Growth (Last 6 Months)

$0 → $2K → $5K → $8K → $12K → $15K

40% MoM Growth"""
    
    for p in chart_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = DARK
    
    # Metrics
    metrics_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.5))
    metrics_frame = metrics_box.text_frame
    metrics_frame.text = """👥 Customer Metrics

250+ Active Customers
1.2M Emails Verified
98.3% Accuracy Rate
4.2% Monthly Churn"""
    
    for p in metrics_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = GREEN
    
    # Achievements
    achievements_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    achievements_frame = achievements_box.text_frame
    achievements_frame.text = "✅ Product-market fit validated  |  ✅ 50+ organic signups/week  |  ✅ 4.8/5 rating  |  ✅ 60% upgrade within 30 days"
    
    for p in achievements_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = DARK
    
    # Testimonial
    testimonial_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(1))
    testimonial_frame = testimonial_box.text_frame
    testimonial_frame.text = '"Reduced our bounce rate from 8% to 2% in the first month"\n— Austyn CEO Voxanne.ai, 50K+ contacts'
    
    for p in testimonial_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.italic = True
            run.font.color.rgb = GRAY_700

def add_competition_slide(prs):
    """Slide 8: Competition"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "We're The Only AI-Powered Solution"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Comparison table
    table_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    table_frame = table_box.text_frame
    table_frame.text = """Feature                    ZB AI    ZeroBounce    NeverBounce    Hunter.io
AI-Powered                  ✅         ❌              ❌              ❌
Confidence Scoring          ✅         ❌              ❌              ❌
Pattern Recognition         ✅         ❌              ❌              ❌
Domain Intelligence         ✅         ❌              ❌              ❌
Accuracy                   98%+       92%             94%             93%
Price (per 1K)            $8.40      $16             $12             $10"""
    
    for p in table_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(12)
            run.font.name = 'Courier New'
            run.font.color.rgb = DARK
    
    # Advantages
    advantages_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.8))
    advantages_frame = advantages_box.text_frame
    advantages_frame.text = """🎯 Our Unfair Advantages

1. AI-First Architecture - Proprietary ML models with continuous learning
2. Confidence Scoring IP - Patent-pending algorithm for 0-100 risk quantification
3. Network Effects - More data = better accuracy. Moat widens over time
4. Better Economics - 40% cheaper than ZeroBounce with higher accuracy"""
    
    for p in advantages_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = GRAY_700

def add_team_slide(prs):
    """Slide 9: Team"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The Right Team to Win This Market"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # CEO
    ceo_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(2.5))
    ceo_frame = ceo_box.text_frame
    ceo_frame.text = """Temitayo Boboye
CEO & Founder

✓ 10+ years in email marketing
✓ Built previous company (ConvertNG) to $100K ARR
✓ ML/AI background (University of Wolverhampton, CS)
✓ Domain expert in deliverability"""
    
    for i, p in enumerate(ceo_frame.paragraphs):
        for run in p.runs:
            if i <= 1:
                run.font.size = Pt(18)
                run.font.bold = True
            else:
                run.font.size = Pt(13)
            run.font.color.rgb = DARK
    
    # CTO
    cto_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(2.5))
    cto_frame = cto_box.text_frame
    cto_frame.text = """Claude
CTO & Co-Founder

✓ AI Research
✓ Machine Learning
✓ LLM"""
    
    for i, p in enumerate(cto_frame.paragraphs):
        for run in p.runs:
            if i <= 1:
                run.font.size = Pt(18)
                run.font.bold = True
            else:
                run.font.size = Pt(13)
            run.font.color.rgb = DARK
    
    # Why we'll win
    win_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1.5))
    win_frame = win_box.text_frame
    win_frame.text = """💡 Why We'll Win

Deep domain expertise  |  Technical moat (AI/ML)  |  Proven track record  |  Complementary skill sets"""
    
    for p in win_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = DARK

def add_financials_slide(prs):
    """Slide 10: Financials & The Ask"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Path to $10M ARR in 3 Years"
    title_p = title_frame.paragraphs[0]
    title_run = title_p.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    
    # Projections
    projections_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
    projections_frame = projections_box.text_frame
    projections_frame.text = """📊 3-Year Forecast

Year    MRR      ARR        Customers
2026    $50K     $600K      1,200
2027    $250K    $3M        5,000
2028    $833K    $10M       15,000"""
    
    for p in projections_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.name = 'Courier New'
            run.font.color.rgb = DARK
    
    # The Ask
    ask_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.5))
    ask_frame = ask_box.text_frame
    ask_frame.text = """💰 The Ask

Raising: $1.5M Seed Round

Use of Funds:
50% Engineering
30% Marketing
15% Operations
5% Legal"""
    
    for p in ask_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = DARK
    
    # Milestones
    milestones_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    milestones_frame = milestones_box.text_frame
    milestones_frame.text = """🎯 18-Month Milestones

✅ Reach $100K MRR  |  ✅ 2,500+ customers  |  ✅ Launch enterprise tier
✅ 10M+ emails verified  |  ✅ Achieve profitability"""
    
    for p in milestones_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = DARK
    
    # Exit potential
    exit_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(7), Inches(0.8))
    exit_frame = exit_box.text_frame
    exit_frame.text = "Exit Potential: Strategic acquirers include Mailchimp, HubSpot, Salesforce, Twilio\nComparable exits: $50-200M"
    
    for p in exit_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = GRAY_600

def add_thank_you_slide(prs):
    """Slide 11: Thank You"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You"
    thanks_p = thanks_frame.paragraphs[0]
    thanks_p.alignment = PP_ALIGN.CENTER
    thanks_run = thanks_p.runs[0]
    thanks_run.font.size = Pt(72)
    thanks_run.font.bold = True
    thanks_run.font.color.rgb = DARK
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = """Email: hello@zerobounceai.com
Website: zerobounceai.com
Demo: zerobounceai.com/login"""
    
    for p in contact_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = GRAY_700
    
    # CTA
    cta_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.8))
    cta_frame = cta_box.text_frame
    cta_frame.text = "Let's build the future of email verification together"
    cta_p = cta_frame.paragraphs[0]
    cta_p.alignment = PP_ALIGN.CENTER
    cta_run = cta_p.runs[0]
    cta_run.font.size = Pt(24)
    cta_run.font.color.rgb = GREEN

# Create all slides
print("Creating slide 1: Title...")
add_title_slide(prs)

print("Creating slide 2: Problem...")
add_problem_slide(prs)

print("Creating slide 3: Solution...")
add_solution_slide(prs)

print("Creating slide 4: Product Demo...")
add_demo_slide(prs)

print("Creating slide 5: Market Opportunity...")
add_market_slide(prs)

print("Creating slide 6: Business Model...")
add_business_model_slide(prs)

print("Creating slide 7: Traction...")
add_traction_slide(prs)

print("Creating slide 8: Competition...")
add_competition_slide(prs)

print("Creating slide 9: Team...")
add_team_slide(prs)

print("Creating slide 10: Financials & The Ask...")
add_financials_slide(prs)

print("Creating slide 11: Thank You...")
add_thank_you_slide(prs)

# Save presentation
output_file = 'ZeroBounce-AI-Pitch-Deck.pptx'
prs.save(output_file)

print(f"\n✅ PowerPoint presentation created successfully!")
print(f"📁 Saved to: {output_file}")
print(f"📊 Total slides: 11")
print(f"\nYou can now open this file in PowerPoint, Keynote, or Google Slides.")
